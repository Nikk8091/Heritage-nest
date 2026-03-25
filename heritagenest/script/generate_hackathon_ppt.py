from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import requests
import json
from io import BytesIO
from urllib.parse import quote
import os
from dotenv import load_dotenv

# Load environment variables
load_dotenv(Path(__file__).resolve().parents[1] / ".env.local")

ROOT = Path(__file__).resolve().parents[1]
OUT_PPT = ROOT / "HeritageNest_Hackathon_Presentation.pptx"
TEMP_PPT = ROOT / "HeritageNest_Hackathon_Presentation_Temp.pptx"
ASSETS_DIR = ROOT / "script" / "ppt_assets"
ASSETS_DIR.mkdir(parents=True, exist_ok=True)

# Firebase & Supabase config
FIREBASE_API_KEY = os.getenv("NEXT_PUBLIC_FIREBASE_API_KEY")
FIREBASE_PROJECT_ID = os.getenv("NEXT_PUBLIC_FIREBASE_PROJECT_ID")
SUPABASE_URL = os.getenv("NEXT_PUBLIC_SUPABASE_URL")

FIRESTORE_REST_URL = f"https://firestore.googleapis.com/v1/projects/{FIREBASE_PROJECT_ID}/databases/default/documents/artItems"

TITLE_COLOR = RGBColor(44, 24, 16)
ACCENT = RGBColor(123, 28, 28)
MUTED = RGBColor(92, 61, 46)
BG = RGBColor(250, 246, 240)

CATEGORIES = [
    ("Folk Art", "Madhubani, Warli, Gond"),
    ("Dance", "Bharatanatyam, Kathak"),
    ("Heritage Site", "Temples, forts, ghats"),
    ("Craft", "Pottery, Dhokra, weaving"),
    ("Festival", "Rituals and celebrations"),
    ("Music", "Folk and classical audio"),
    ("Textile", "Phulkari, Kalamkari"),
    ("Sculpture", "Stone, terracotta, metal"),
]

TEAM_MEMBERS = [
    "Nikhil Kumar",
    "Lachydeep",
    "Sweta Sharma",
    "Sneha Sharma",
    "Aastha",
]


def fetch_items_by_category(category):
    """Fetch art items from Firestore for a specific category"""
    try:
        # Build Firestore REST query
        query_params = {
            "pageSize": 6,
            "structuredQuery": {
                "from": [{"collectionId": "artItems"}],
                "where": {
                    "compositeFilter": {
                        "op": "AND",
                        "filters": [
                            {
                                "fieldFilter": {
                                    "field": {"fieldPath": "category"},
                                    "op": "EQUAL",
                                    "value": {"stringValue": category}
                                }
                            },
                            {
                                "fieldFilter": {
                                    "field": {"fieldPath": "moderation_status"},
                                    "op": "EQUAL",
                                    "value": {"stringValue": "approved"}
                                }
                            }
                        ]
                    }
                },
                "orderBy": [{"field": {"fieldPath": "created_at"}, "direction": "DESCENDING"}]
            }
        }
        
        params = {"key": FIREBASE_API_KEY}
        response = requests.post(
            f"{FIRESTORE_REST_URL}:runQuery",
            json=query_params,
            params=params,
            timeout=10
        )
        
        items = []
        if response.status_code == 200:
            results = response.json()
            for result in results:
                if "document" in result:
                    doc_data = result["document"].get("fields", {})
                    # Parse Firestore field format
                    item = {
                        "id": result["document"]["name"].split("/")[-1],
                        "title": doc_data.get("title", {}).get("stringValue", ""),
                        "description": doc_data.get("description", {}).get("stringValue", ""),
                        "media_url": doc_data.get("media_url", {}).get("stringValue", ""),
                        "art_form": doc_data.get("art_form", {}).get("stringValue", ""),
                        "state": doc_data.get("state", {}).get("stringValue", ""),
                    }
                    if item["media_url"]:
                        items.append(item)
        
        return items[:4]  # Return up to 4 items
    except Exception as e:
        print(f"⚠️  Error fetching items for {category}: {e}")
        return []


def download_image(url, max_width=400, max_height=300):
    """Download image from URL and return PIL Image object"""
    try:
        if not url:
            return None
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        img = Image.open(BytesIO(response.content)).convert("RGB")
        
        # Resize to fit in bounds while maintaining aspect ratio
        img.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
        return img
    except Exception as e:
        print(f"⚠️  Error downloading image from {url}: {e}")
        return None


def create_category_collage(category_name, items, output_path):
    """Create a collage of images for a category"""
    try:
        # Create canvas
        canvas = Image.new("RGB", (1280, 720), (250, 246, 240))
        
        # Download images
        images = []
        for item in items[:4]:
            img = download_image(item.get("media_url"))
            if img:
                images.append((img, item.get("title", "Untitled")))
        
        if not images:
            # Fallback to gradient if no images
            return make_category_card(output_path, category_name, f"{len(items)} items", (123, 28, 28), (193, 96, 58))
        
        # Arrange images in grid
        x_positions = [50, 340, 630]
        y_position = 100
        
        for idx, (img, title) in enumerate(images[:3]):
            # Resize to fit grid
            img.thumbnail((250, 280), Image.Resampling.LANCZOS)
            canvas.paste(img, (x_positions[idx], y_position))
        
        # Add category title and info
        draw = ImageDraw.Draw(canvas)
        try:
            font_title = ImageFont.truetype("arial.ttf", 56)
            font_info = ImageFont.truetype("arial.ttf", 24)
        except:
            font_title = ImageFont.load_default()
            font_info = ImageFont.load_default()
        
        # Semi-transparent background for text
        draw.rectangle([50, 420, 1230, 670], fill=(255, 255, 255, 220))
        
        draw.text((70, 440), category_name, font=font_title, fill=(44, 24, 16))
        draw.text((70, 510), f"{len(items)} items in database", font=font_info, fill=(92, 61, 46))
        draw.text((70, 550), ", ".join([item.get("art_form", "")[:20] for item in items[:3] if item.get("art_form")]), 
                  font=font_info, fill=(123, 28, 28))
        
        canvas.save(output_path)
        return
    except Exception as e:
        print(f"⚠️  Error creating collage: {e}")
        return make_category_card(output_path, category_name, f"{len(items)} items", (123, 28, 28), (193, 96, 58))


def make_category_card(path: Path, title: str, subtitle: str, color_a, color_b):
    """Create a gradient card as fallback when no images are available"""
    img = Image.new("RGB", (1280, 720), color_a)
    draw = ImageDraw.Draw(img)

    # Gradient background
    for y in range(720):
        blend = y / 719
        r = int(color_a[0] * (1 - blend) + color_b[0] * blend)
        g = int(color_a[1] * (1 - blend) + color_b[1] * blend)
        b = int(color_a[2] * (1 - blend) + color_b[2] * blend)
        draw.line([(0, y), (1280, y)], fill=(r, g, b))

    # Decorative pattern
    for i in range(0, 1280, 80):
        draw.rectangle([i, 0, i + 40, 720], outline=(255, 255, 255, 22), width=1)

    try:
        font_title = ImageFont.truetype("arial.ttf", 72)
        font_sub = ImageFont.truetype("arial.ttf", 36)
    except Exception:
        font_title = ImageFont.load_default()
        font_sub = ImageFont.load_default()

    # Content box
    draw.rounded_rectangle([90, 130, 1190, 590], radius=40, fill=(255, 255, 255, 210))
    draw.text((150, 230), title, font=font_title, fill=(50, 30, 20))
    draw.text((150, 340), subtitle, font=font_sub, fill=(92, 61, 46))
    draw.text((150, 450), "HeritageNest Category Showcase", font=font_sub, fill=(123, 28, 28))

    img.save(path)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(12.0), Inches(1.0))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.82), Inches(1.35), Inches(11.5), Inches(0.8))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = MUTED


def add_bullets(slide, items, left=0.9, top=2.0, width=11.4, height=4.8, size=22):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(45, 35, 30)


def add_footer(slide, text="Hackathon 2026 | HeritageNest"):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(12.0), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(120, 95, 78)


def build_presentation():
    print("📊 Fetching items from Firebase Firestore...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Fetch category data and create card images
    card_images = []
    category_stats = []
    
    palettes = [
        ((123, 28, 28), (193, 96, 58)),
        ((44, 62, 80), (52, 73, 94)),
        ((26, 74, 114), (43, 117, 162)),
        ((80, 40, 20), (140, 92, 58)),
        ((102, 51, 0), (168, 112, 40)),
        ((45, 76, 62), (84, 126, 97)),
        ((105, 46, 74), (159, 79, 114)),
        ((58, 72, 92), (98, 113, 136)),
    ]

    for idx, (category, subtitle) in enumerate(CATEGORIES):
        print(f"  📸 Processing {category}...")
        image_path = ASSETS_DIR / f"category_{idx+1}.png"
        
        # Fetch real items
        items = fetch_items_by_category(category)
        category_stats.append((category, len(items)))
        
        if items:
            # Create collage with real images
            create_category_collage(category, items, image_path)
        else:
            # Fallback to gradient card
            a, b = palettes[idx % len(palettes)]
            make_category_card(image_path, category, subtitle, a, b)
        
        card_images.append(image_path)

    # Slide 1: Title
    print("📝 Creating presentation slides...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "HeritageNest", "Regional Folk Art Digital Archive | Hackathon Presentation")
    add_bullets(slide, [
        "A community-driven platform preserving India's intangible cultural heritage.",
        "Supports image, video, shorts, and audio artifacts with moderation and AI explanations.",
        "Built for creators, learners, and cultural communities.",
    ], top=2.1, size=24)
    add_bullets(
        slide,
        ["Team: " + ", ".join(TEAM_MEMBERS)],
        top=5.8,
        size=16,
    )
    add_footer(slide, "Team Contact: Overall Manager nikk4645@gmail.com | Tech Support swetashr08@gmail.com")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Problem Statement", "Why this project matters")
    add_bullets(slide, [
        "Traditional knowledge is scattered and often undocumented in digital form.",
        "Audio/video folk practices are hard to discover and preserve in one place.",
        "Youth engagement needs modern, searchable, social-first formats.",
    ], size=23)
    add_footer(slide)

    # New Slide: Database Statistics
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Current Database Statistics", "Real data from our live platform")
    stats_text = []
    total_items = sum(count for _, count in category_stats)
    for cat, count in category_stats:
        stats_text.append(f"• {cat}: {count} items")
    stats_text.insert(0, f"📊 Total Items in Archive: {total_items}")
    add_bullets(slide, stats_text, top=2.0, size=20)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Solution Overview", "What HeritageNest delivers")
    add_bullets(slide, [
        "Unified archive: Artcards + Shorts + Audio heritage library.",
        "Category filtering across Folk Art, Dance, Music, Craft, Textile, and more.",
        "Admin moderation workflow with pending/approved/rejected pipeline.",
        "Role-aware access: public audience, contributors, and admins.",
    ], size=22)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "System Architecture", "Modern full-stack build")
    add_bullets(slide, [
        "Frontend: Next.js App Router + React + CSS Modules",
        "Auth + Database: Firebase Auth + Firestore",
        "Media Storage: Supabase Storage (public bucket)",
        "AI Layer: FastAPI endpoint + Groq Llama fallback",
        "Deploy: Vercel (web) with managed environment variables",
    ], size=20)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Category Artcard Showcase I", "Different category visual cards")
    slide.shapes.add_picture(str(card_images[0]), Inches(0.8), Inches(1.7), width=Inches(4.0))
    slide.shapes.add_picture(str(card_images[1]), Inches(4.7), Inches(1.7), width=Inches(4.0))
    slide.shapes.add_picture(str(card_images[2]), Inches(8.6), Inches(1.7), width=Inches(4.0))
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Category Artcard Showcase II", "Different category visual cards")
    slide.shapes.add_picture(str(card_images[3]), Inches(0.8), Inches(1.7), width=Inches(4.0))
    slide.shapes.add_picture(str(card_images[4]), Inches(4.7), Inches(1.7), width=Inches(4.0))
    slide.shapes.add_picture(str(card_images[5]), Inches(8.6), Inches(1.7), width=Inches(4.0))
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Category Artcard Showcase III", "Textile and Sculpture categories")
    slide.shapes.add_picture(str(card_images[6]), Inches(1.8), Inches(1.7), width=Inches(4.5))
    slide.shapes.add_picture(str(card_images[7]), Inches(7.0), Inches(1.7), width=Inches(4.5))
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Feature Walkthrough", "User journeys in one platform")
    add_bullets(slide, [
        "Creator upload modes: Artcard, Short Video, and Audio with metadata.",
        "Discovery: Search, filters, category pages, and dedicated Audio section.",
        "Engagement: Save/bookmark, share, and AI learning companion.",
        "Admin tools: moderation queue, item management, and insights dashboard.",
    ], size=21)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "AI Model Combination", "Hybrid intelligence approach")
    add_bullets(slide, [
        "Primary model path: FastAPI cultural assistant service.",
        "Fallback model path: Groq Llama-3.1-8b-instant for resilience.",
        "Rule model layer: scope guardrails for topic compliance and safe refusal.",
        "Context model layer: item metadata grounding to reduce hallucinations.",
        "Result: Reliable, domain-focused explanations for users.",
    ], size=20)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Hackathon Impact & Next Steps", "Roadmap after MVP")
    add_bullets(slide, [
        "Impact: stronger digital preservation and easier cultural discovery.",
        "Next: multilingual assistant + regional script support.",
        "Next: richer audio player analytics and creator performance insights.",
        "Next: community verification and partner museum integrations.",
    ], size=21)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Thank You", "Demo + Q&A")
    add_bullets(slide, [
        "Live App: https://heritage-nest-rtm6.vercel.app",
        "Team: " + ", ".join(TEAM_MEMBERS),
        "Overall Manager: nikk4645@gmail.com",
        "Tech Support: swetashr08@gmail.com",
    ], top=2.4, size=24)
    add_footer(slide, "HeritageNest | Preserving culture through technology")

    # Save to temp file then rename
    prs.save(TEMP_PPT)
    
    # Try to move/rename
    import shutil
    import time
    try:
        # Wait a moment for file to be fully written
        time.sleep(0.5)
        # Remove old if exists
        if OUT_PPT.exists():
            try:
                OUT_PPT.unlink(missing_ok=True)
            except:
                pass
        # Move temp to final
        shutil.move(str(TEMP_PPT), str(OUT_PPT))
    except Exception as e:
        # If all else fails, just copy it
        try:
            shutil.copy2(TEMP_PPT, OUT_PPT)
            TEMP_PPT.unlink(missing_ok=True)
        except Exception as copy_e:
            print(f"⚠️  Could not rename/copy: {e}, {copy_e}")
            print(f"📍 Presentation saved as: {TEMP_PPT}")


if __name__ == "__main__":
    build_presentation()
    print(f"✅ Generated: {OUT_PPT}")
